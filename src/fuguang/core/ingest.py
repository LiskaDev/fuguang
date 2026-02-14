# ingest.py - 知识吞噬系统 (Knowledge Eater)
"""
将本地文件（PDF/Word/TXT/代码）导入向量数据库

功能：
- 支持多种格式：PDF, DOCX, TXT, MD, PY, JSON, LOG, CSV
- 智能分块：按段落/句子边界切分，避免上下文断裂
- 批量导入：支持整个文件夹
- 进度追踪：显示导入进度
"""

import os
import re
import logging
from typing import Optional
from pathlib import Path

logger = logging.getLogger("Fuguang")

# 尝试导入文件解析库
try:
    from pypdf import PdfReader
    PYPDF_AVAILABLE = True
except ImportError:
    PYPDF_AVAILABLE = False
    logger.warning("⚠️ pypdf 未安装，PDF 导入功能将不可用")

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("⚠️ python-docx 未安装，Word 导入功能将不可用")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("⚠️ python-pptx 未安装，PowerPoint 导入功能将不可用")

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    logger.warning("⚠️ openpyxl 未安装，Excel 导入功能将不可用")


class KnowledgeEater:
    """知识吞噬者 - 将文件内容导入向量数据库"""
    
    # 支持的文件格式
    SUPPORTED_EXTENSIONS = {
        'pdf': '📕 PDF文档',
        'docx': '📘 Word文档',
        'pptx': '📊 PowerPoint演示',
        'xlsx': '📗 Excel表格',
        'xls': '📗 Excel表格(旧版)',
        'txt': '📄 纯文本',
        'md': '📝 Markdown',
        'py': '🐍 Python代码',
        'js': '📜 JavaScript',
        'json': '📋 JSON数据',
        'log': '📃 日志文件',
        'csv': '📊 CSV表格',
        'html': '🌐 HTML网页',
        'xml': '📰 XML文件',
    }
    
    def __init__(self, memory_bank, chunk_size: int = 500, overlap: int = 50):
        """
        初始化知识吞噬者
        
        Args:
            memory_bank: MemoryBank 实例
            chunk_size: 每个文本块的目标字符数
            overlap: 相邻块之间的重叠字符数
        """
        self.memory = memory_bank
        self.chunk_size = chunk_size
        self.overlap = overlap
        
    def ingest_file(self, file_path: str) -> str:
        """
        吞噬单个文件
        
        Args:
            file_path: 文件的绝对路径
            
        Returns:
            操作结果消息
        """
        path = Path(file_path)
        
        # 检查文件是否存在
        if not path.exists():
            return f"❌ 文件不存在: {file_path}"
            
        if not path.is_file():
            return f"❌ 这不是文件: {file_path}"
        
        ext = path.suffix.lower().lstrip('.')
        
        # 检查格式支持
        if ext not in self.SUPPORTED_EXTENSIONS:
            supported = ', '.join(self.SUPPORTED_EXTENSIONS.keys())
            return f"❌ 不支持的格式 .{ext}。支持: {supported}"
        
        logger.info(f"🍽️ [进食] 正在读取: {path.name}")
        
        # 提取文本
        try:
            text = self._extract_text(path, ext)
        except Exception as e:
            logger.error(f"❌ 提取文本失败: {e}")
            return f"❌ 读取失败: {str(e)}"
        
        if not text or not text.strip():
            return f"⚠️ 文件是空的或无法提取文本: {path.name}"
        
        # 分块
        chunks = self._smart_chunk(text)
        logger.info(f"🔪 [消化] 切分为 {len(chunks)} 个碎片")
        
        # 存入向量数据库
        success_count = 0
        for i, chunk in enumerate(chunks):
            try:
                metadata = {
                    "source": path.name,
                    "source_path": str(path),
                    "type": "document",
                    "format": ext,
                    "chunk_id": i,
                    "total_chunks": len(chunks)
                }
                self.memory.add_memory(chunk, category="knowledge", metadata=metadata)
                success_count += 1
                
                # 进度提示（每 20 个块显示一次）
                if (i + 1) % 20 == 0:
                    progress = int((i + 1) / len(chunks) * 100)
                    logger.info(f"   📦 进度: {progress}% ({i + 1}/{len(chunks)})")
                    
            except Exception as e:
                logger.error(f"   ❌ 第 {i} 块存储失败: {e}")
        
        result = f"✅ 已吞噬 '{path.name}'，存入 {success_count}/{len(chunks)} 条知识碎片"
        logger.info(f"🎉 {result}")
        return result
    
    def ingest_folder(self, folder_path: str, recursive: bool = True) -> str:
        """
        批量吞噬文件夹中的所有文件
        
        Args:
            folder_path: 文件夹路径
            recursive: 是否递归处理子文件夹
            
        Returns:
            操作结果消息
        """
        folder = Path(folder_path)
        
        if not folder.exists():
            return f"❌ 文件夹不存在: {folder_path}"
            
        if not folder.is_dir():
            return f"❌ 这不是文件夹: {folder_path}"
        
        # 收集所有支持的文件
        pattern = "**/*" if recursive else "*"
        files = []
        for ext in self.SUPPORTED_EXTENSIONS.keys():
            files.extend(folder.glob(f"{pattern}.{ext}"))
        
        if not files:
            return f"⚠️ 文件夹中没有找到支持的文件"
        
        logger.info(f"📂 [批量进食] 发现 {len(files)} 个文件")
        
        results = []
        for i, file_path in enumerate(files):
            logger.info(f"📖 [{i+1}/{len(files)}] {file_path.name}")
            result = self.ingest_file(str(file_path))
            results.append(f"{file_path.name}: {result}")
        
        return f"📚 批量吞噬完成！处理了 {len(files)} 个文件"
    
    def _extract_text(self, path: Path, ext: str) -> str:
        """根据文件格式提取文本"""
        
        if ext == 'pdf':
            if not PYPDF_AVAILABLE:
                raise RuntimeError("pypdf 未安装，请运行: pip install pypdf")
            reader = PdfReader(str(path))
            text_parts = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            return "\n\n".join(text_parts)
        
        elif ext == 'docx':
            if not DOCX_AVAILABLE:
                raise RuntimeError("python-docx 未安装，请运行: pip install python-docx")
            doc = docx.Document(str(path))
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            return "\n\n".join(paragraphs)
        
        elif ext == 'pptx':
            if not PPTX_AVAILABLE:
                raise RuntimeError("python-pptx 未安装，请运行: pip install python-pptx")
            prs = Presentation(str(path))
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    text_parts.append(f"[幻灯片{slide_num}] " + "\n".join(slide_texts))
            return "\n\n".join(text_parts)
        
        elif ext in ('xlsx', 'xls'):
            if not OPENPYXL_AVAILABLE:
                raise RuntimeError("openpyxl 未安装，请运行: pip install openpyxl")
            wb = openpyxl.load_workbook(str(path), data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else '' for c in row]
                    if any(c.strip() for c in cells):
                        rows.append(' | '.join(cells))
                if rows:
                    text_parts.append(f"[工作表: {sheet_name}]\n" + "\n".join(rows))
            return "\n\n".join(text_parts)
        
        else:
            # 纯文本类文件
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            for encoding in encodings:
                try:
                    with open(path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            # 最后尝试忽略错误
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
    
    def _smart_chunk(self, text: str) -> list:
        """
        智能分块：尽量在句子/段落边界切分
        
        优先级：
        1. 双换行（段落边界）
        2. 单换行
        3. 句号/问号/感叹号
        4. 逗号/分号
        5. 空格
        6. 强制按字符数切分
        """
        chunks = []
        current_chunk = ""
        
        # 按段落初步切分
        paragraphs = re.split(r'\n\s*\n', text)
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
                
            # 如果段落本身就很长，需要进一步切分
            if len(para) > self.chunk_size:
                # 按句子切分
                sentences = re.split(r'([。！？!?])', para)
                temp = ""
                for i, part in enumerate(sentences):
                    temp += part
                    # 如果是标点符号后面，或者积累够长了
                    if len(temp) >= self.chunk_size:
                        if current_chunk:
                            chunks.append(current_chunk)
                        current_chunk = temp[-self.overlap:] if len(temp) > self.overlap else temp
                        temp = temp[:-self.overlap] if len(temp) > self.overlap else ""
                        chunks.append(temp)
                        temp = current_chunk
                        current_chunk = ""
                if temp:
                    if len(current_chunk) + len(temp) <= self.chunk_size:
                        current_chunk += temp
                    else:
                        if current_chunk:
                            chunks.append(current_chunk)
                        current_chunk = temp
            else:
                # 段落较短，尝试合并
                if len(current_chunk) + len(para) + 2 <= self.chunk_size:
                    current_chunk += ("\n\n" + para if current_chunk else para)
                else:
                    if current_chunk:
                        chunks.append(current_chunk)
                    current_chunk = para
        
        # 最后一个块
        if current_chunk:
            chunks.append(current_chunk)
        
        # 过滤太短的块（少于 50 字符）
        chunks = [c for c in chunks if len(c.strip()) >= 10]
        
        return chunks
    
    def get_supported_formats(self) -> str:
        """获取支持的文件格式列表"""
        lines = ["支持的文件格式:"]
        for ext, desc in self.SUPPORTED_EXTENSIONS.items():
            lines.append(f"  .{ext} - {desc}")
        return "\n".join(lines)
