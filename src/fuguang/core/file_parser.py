# ==================================================
# 📄 file_parser.py - 扶光文件解析公共模块
# ==================================================
# 从 qq_bridge.py 抽取，供 QQBridge 和 WebBridge 共用
# ==================================================

import os
import tempfile
import logging
from typing import Optional, Callable

logger = logging.getLogger("Fuguang")

# 支持的扩展名
SUPPORTED_EXTENSIONS = {".txt", ".pdf", ".docx", ".pptx", ".xlsx", ".jpg", ".jpeg", ".png"}
MAX_CONTENT_LENGTH = 3000


def parse_file(file_path: str, file_name: str,
               image_analyzer: Optional[Callable[[str], str]] = None) -> str:
    """
    根据扩展名解析文件内容（同步方法）

    Args:
        file_path:       文件的本地路径
        file_name:       原始文件名（用于判断扩展名）
        image_analyzer:  图片分析回调 fn(image_path) -> description
                         如不提供，图片相关内容标记为 [图片：未配置视觉分析]

    Returns:
        解析后的文本内容（截断至 MAX_CONTENT_LENGTH 字符）
    """
    ext = os.path.splitext(file_name)[1].lower()
    try:
        if ext == ".txt":
            return parse_txt(file_path)
        elif ext == ".pdf":
            return parse_pdf(file_path, image_analyzer)
        elif ext == ".docx":
            return parse_docx(file_path, image_analyzer)
        elif ext == ".pptx":
            return parse_pptx(file_path, image_analyzer)
        elif ext == ".xlsx":
            return parse_xlsx(file_path)
        elif ext in (".jpg", ".jpeg", ".png"):
            if image_analyzer:
                return image_analyzer(file_path)
            return "[图片：未配置视觉分析]"
        else:
            return "暂不支持此格式"
    except Exception as e:
        logger.error(f"📄 文件解析失败 ({file_name}): {e}")
        return f"文件解析失败: {e}"


def _analyze_image(image_path: str,
                   image_analyzer: Optional[Callable[[str], str]]) -> str:
    """调用图片分析回调，未配置时返回占位文本"""
    if image_analyzer:
        try:
            return image_analyzer(image_path)
        except Exception as e:
            return f"(图片分析失败: {e})"
    return "[图片：未配置视觉分析]"


# ==================================================
# 各格式解析函数
# ==================================================

def parse_txt(file_path: str) -> str:
    """解析 TXT 文件"""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    return content[:MAX_CONTENT_LENGTH]


def parse_pdf(file_path: str,
              image_analyzer: Optional[Callable[[str], str]] = None) -> str:
    """解析 PDF 文件：逐页提取文字，文字不足 20 字的页面截图分析"""
    import fitz  # pymupdf
    result_parts = []
    tmp_files = []
    try:
        doc = fitz.open(file_path)
        for i, page in enumerate(doc):
            page_text = page.get_text().strip()
            if len(page_text) >= 20:
                result_parts.append(f"--- 第{i+1}页 ---\n{page_text}")
            else:
                # 文字太少，可能是扫描件/图片，截图分析
                try:
                    pix = page.get_pixmap(dpi=150)
                    tmp_img = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
                    tmp_img.close()
                    pix.save(tmp_img.name)
                    tmp_files.append(tmp_img.name)
                    img_desc = _analyze_image(tmp_img.name, image_analyzer)
                    result_parts.append(f"--- 第{i+1}页(图片) ---\n{img_desc}")
                except Exception:
                    result_parts.append(f"--- 第{i+1}页 ---\n{page_text or '(无法识别)'}")
        doc.close()
    finally:
        for tmp in tmp_files:
            try:
                os.unlink(tmp)
            except Exception:
                pass
    content = "\n".join(result_parts)
    return content[:MAX_CONTENT_LENGTH]


def parse_docx(file_path: str,
               image_analyzer: Optional[Callable[[str], str]] = None) -> str:
    """解析 DOCX 文件：提取段落文字 + inline_shapes 图片"""
    from docx import Document
    result_parts = []
    tmp_files = []
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                result_parts.append(para.text.strip())
        # 提取嵌入图片
        for shape in doc.inline_shapes:
            try:
                blip = shape._inline.graphic.graphicData.pic.blipFill.blip
                rId = blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                if rId:
                    image_part = doc.part.related_parts[rId]
                    tmp_img = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
                    tmp_img.write(image_part.blob)
                    tmp_img.close()
                    tmp_files.append(tmp_img.name)
                    img_desc = _analyze_image(tmp_img.name, image_analyzer)
                    result_parts.append(f"[文档内图片：{img_desc}]")
            except Exception:
                result_parts.append("[文档内图片：无法提取]")
    finally:
        for tmp in tmp_files:
            try:
                os.unlink(tmp)
            except Exception:
                pass
    content = "\n".join(result_parts)
    return content[:MAX_CONTENT_LENGTH]


def parse_pptx(file_path: str,
               image_analyzer: Optional[Callable[[str], str]] = None) -> str:
    """解析 PPTX 文件：逐页提取文字框 + Picture 图片"""
    from pptx import Presentation
    from pptx.shapes.picture import Picture
    result_parts = []
    tmp_files = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if isinstance(shape, Picture):
                    try:
                        img_blob = shape.image.blob
                        tmp_img = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
                        tmp_img.write(img_blob)
                        tmp_img.close()
                        tmp_files.append(tmp_img.name)
                        img_desc = _analyze_image(tmp_img.name, image_analyzer)
                        slide_texts.append(f"[幻灯片图片：{img_desc}]")
                    except Exception:
                        slide_texts.append("[幻灯片图片：无法提取]")
                elif shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text.strip())
            if slide_texts:
                result_parts.append(f"--- 第{i+1}页 ---\n" + "\n".join(slide_texts))
    finally:
        for tmp in tmp_files:
            try:
                os.unlink(tmp)
            except Exception:
                pass
    content = "\n".join(result_parts)
    return content[:MAX_CONTENT_LENGTH]


def parse_xlsx(file_path: str) -> str:
    """解析 XLSX 文件：读取所有 sheet 的单元格内容"""
    from openpyxl import load_workbook
    result_parts = []
    wb = load_workbook(file_path, read_only=True, data_only=True)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        result_parts.append(f"{sheet_name}:")
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None:
                    result_parts.append(f"第{cell.row}行第{cell.column}列: {cell.value}")
    wb.close()
    content = "\n".join(result_parts)
    return content[:MAX_CONTENT_LENGTH]
